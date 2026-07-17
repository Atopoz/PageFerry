"""覆盖 PPTX 结构保真翻译的 synthetic regression tests."""

from collections.abc import Callable, Sequence
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from modules.pptx import PptxPipeline, PptxRunNormalizer
from modules.pptx.markup import deterministic_repair
from modules.translation.contracts import (
    TranslationBatchItem,
    TranslationBatchResult,
    TranslationProgress,
)


class RecordingTranslator:
    """使用确定性变换生成译文, 并记录每个 format hint."""

    def __init__(self, transform: Callable[[str], str]) -> None:
        """保存每个 batch item 共用的文本变换函数."""

        self._transform = transform
        self.calls: list[tuple[str, tuple[str, ...]]] = []

    def translate_batch(
        self,
        *,
        texts: Sequence[str],
        source_language: str | None,
        target_language: str,
        format_hint: str,
        read_only_context: Sequence[str] = (),
        repair_candidates: Sequence[str] = (),
    ) -> TranslationBatchResult:
        """按 provider adapter 的 contract 返回带索引的变换结果."""

        del source_language, target_language, read_only_context, repair_candidates
        captured = tuple(texts)
        self.calls.append((format_hint, captured))
        return TranslationBatchResult(
            items=tuple(
                TranslationBatchItem(index=index, text=self._transform(text))
                for index, text in enumerate(captured)
            )
        )


class FailingTranslator:
    """每次调用都抛出异常, 让 fallback 行为保持确定."""

    def __init__(self) -> None:
        """记录正常调用, 确保异常 batch 不会触发逐段 repair。"""

        self.format_hints: list[str] = []

    def translate_batch(
        self,
        *,
        texts: Sequence[str],
        source_language: str | None,
        target_language: str,
        format_hint: str,
        read_only_context: Sequence[str] = (),
        repair_candidates: Sequence[str] = (),
    ) -> TranslationBatchResult:
        """记录 hint, 并模拟 provider 故障."""

        del texts, source_language, target_language, read_only_context, repair_candidates
        self.format_hints.append(format_hint)
        raise RuntimeError("synthetic provider failure")


class MissingCandidateTranslator:
    """返回空结果, 模拟 provider 遗漏整个 candidate。"""

    def __init__(self) -> None:
        """初始化 format hint 调用记录。"""

        self.format_hints: list[str] = []

    def translate_batch(
        self,
        *,
        texts: Sequence[str],
        source_language: str | None,
        target_language: str,
        format_hint: str,
        read_only_context: Sequence[str] = (),
        repair_candidates: Sequence[str] = (),
    ) -> TranslationBatchResult:
        """记录调用并返回没有任何 item 的合法 batch 外壳。"""

        del texts, source_language, target_language, read_only_context, repair_candidates
        self.format_hints.append(format_hint)
        return TranslationBatchResult(items=())


def test_deterministic_repair_aligns_missing_spans_and_space_slot() -> None:
    """缺少 wrapper 时按保守规则补齐 span 并恢复 SPACE slot。"""

    source = "<span>前半句</span><span><SPACE></span><span>后半句</span>"

    assert deterministic_repair("Translated sentence.", source) == (
        "<span>Translated sentence.</span><span><SPACE></span><span></span>"
    )


def test_pipeline_translates_shape_group_space_and_speaker_notes(tmp_path: Path) -> None:
    """翻译 shape 和 notes, 同时保留 group 路径, 空格和 notes relationship."""

    source = tmp_path / "source.pptx"
    output = tmp_path / "translated.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    paragraph = text_box.text_frame.paragraphs[0]
    first_run = paragraph.add_run()
    first_run.text = "Hello"
    first_run.font.size = Pt(20)
    whitespace_run = paragraph.add_run()
    whitespace_run.text = " "
    whitespace_run.font.bold = True
    final_run = paragraph.add_run()
    final_run.text = "world"

    group = slide.shapes.add_group_shape()
    child = group.shapes.add_textbox(Inches(1), Inches(3), Inches(4), Inches(1))
    child.text = "Grouped text"

    notes_text_frame = slide.notes_slide.notes_text_frame
    assert notes_text_frame is not None
    notes_text_frame.text = "Private speaker note"
    source_notes_relation = _notes_relationship(slide)
    source_non_body_notes = _non_body_notes_placeholders(slide)
    presentation.save(source)
    original_bytes = source.read_bytes()

    replacements = {
        "Hello": "A substantially longer translated greeting",
        "world": "planet",
        "Grouped text": "Translated grouped text",
        "Private speaker note": "Translated private speaker note",
    }

    def transform(text: str) -> str:
        """只替换可见源文本, 保持所有结构 marker 不变."""

        for source_text, translated_text in replacements.items():
            text = text.replace(source_text, translated_text)
        return text

    translator = RecordingTranslator(transform)
    updates: list[TranslationProgress] = []
    result = PptxPipeline(translator, translate_tables=False).translate_to(
        source_path=source,
        output_path=output,
        source_language="en",
        target_language="fr",
        report_progress=updates.append,
    )

    translated = Presentation(output)
    translated_slide = translated.slides[0]
    translated_paragraph = translated_slide.shapes[0].text_frame.paragraphs[0]
    assert translated_paragraph.text == ("A substantially longer translated greeting planet")
    assert translated_paragraph.runs[1].text == " "
    assert translated_paragraph.runs[0].font.size.pt <= 20
    assert translated_slide.shapes[1].shapes[0].text == "Translated grouped text"
    assert translated_slide.notes_slide.notes_text_frame.text == ("Translated private speaker note")
    assert _notes_relationship(translated_slide) == source_notes_relation
    assert _non_body_notes_placeholders(translated_slide) == source_non_body_notes
    assert source.read_bytes() == original_bytes
    assert result.translated_segments == 3
    assert result.fallback_segments == 0
    assert updates[0] == TranslationProgress(stage="extracting")
    assert updates[1] == TranslationProgress(stage="translating", total_segments=3)
    assert updates[-1] == TranslationProgress(
        stage="formatting",
        processed_segments=3,
        total_segments=3,
    )
    assert {format_hint for format_hint, _ in translator.calls} == {
        "pptx",
        "pptx_notes",
    }
    body_payload = next(texts for hint, texts in translator.calls if hint == "pptx")
    assert any("<span><SPACE></span>" in text for text in body_payload)
    assert any("[SHAPE_2.1]" in text for text in body_payload)


def test_pipeline_translates_table_cells(tmp_path: Path) -> None:
    """翻译 cell 段落, 不重建 table."""

    source = tmp_path / "table-source.pptx"
    output = tmp_path / "table-translated.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table_shape = slide.shapes.add_table(1, 2, Inches(1), Inches(1), Inches(6), Inches(1))
    table_shape.table.cell(0, 0).text = "Alpha"
    table_shape.table.cell(0, 1).text = "Beta"
    presentation.save(source)

    translator = RecordingTranslator(lambda text: text.replace("Alpha", "一").replace("Beta", "二"))
    result = PptxPipeline(translator, translate_notes=False).translate_to(
        source_path=source,
        output_path=output,
        source_language="en",
        target_language="zh-CN",
    )

    translated_table = Presentation(output).slides[0].shapes[0].table
    assert translated_table.cell(0, 0).text == "一"
    assert translated_table.cell(0, 1).text == "二"
    translated_run = translated_table.cell(0, 0).text_frame.paragraphs[0].runs[0]
    run_properties = translated_run._r.get_or_add_rPr()
    assert run_properties.find(qn("a:ea")).get("typeface") == "Noto Sans SC"
    assert run_properties.find(qn("a:latin")) is None
    assert run_properties.find(qn("a:cs")) is None
    assert result.translated_segments == 2
    assert result.fallback_segments == 0
    assert [hint for hint, _ in translator.calls] == ["pptx_table"]


def test_table_skips_numeric_symbol_cell_without_calling_model(tmp_path: Path) -> None:
    """纯数字/符号 cell 保持原样, 且不进入 provider payload。"""

    source = tmp_path / "table-skip-source.pptx"
    output = tmp_path / "table-skip-translated.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table_shape = slide.shapes.add_table(1, 2, Inches(1), Inches(1), Inches(6), Inches(1))
    table_shape.table.cell(0, 0).text = "$99.00 / 2026-07-16"
    table_shape.table.cell(0, 1).text = "Total"
    presentation.save(source)

    translator = RecordingTranslator(lambda text: text.replace("Total", "合计"))
    result = PptxPipeline(translator, translate_notes=False).translate_to(
        source_path=source,
        output_path=output,
        source_language="en",
        target_language="zh-CN",
    )

    translated_table = Presentation(output).slides[0].shapes[0].table
    assert translated_table.cell(0, 0).text == "$99.00 / 2026-07-16"
    assert translated_table.cell(0, 1).text == "合计"
    payload = "".join(text for _, texts in translator.calls for text in texts)
    assert "$99.00" not in payload
    assert "Total" in payload
    assert result.translated_segments == 1


def test_english_translation_restores_spacing_between_run_spans(tmp_path: Path) -> None:
    """English 译文跨视觉 run 时补回单词间空格。"""

    source = tmp_path / "english-spacing-source.pptx"
    output = tmp_path / "english-spacing-translated.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    paragraph = slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(5), Inches(1)
    ).text_frame.paragraphs[0]
    first_run = paragraph.add_run()
    first_run.text = "智能"
    first_run.font.bold = True
    second_run = paragraph.add_run()
    second_run.text = "助手"
    presentation.save(source)

    translator = RecordingTranslator(
        lambda text: text.replace("智能", "Smart").replace("助手", "Assistant")
    )
    PptxPipeline(translator, translate_tables=False, translate_notes=False).translate_to(
        source_path=source,
        output_path=output,
        source_language="zh-CN",
        target_language="en",
    )

    translated_paragraph = Presentation(output).slides[0].shapes[0].text_frame.paragraphs[0]
    assert translated_paragraph.text == "Smart Assistant"
    assert [run.text for run in translated_paragraph.runs] == ["Smart ", "Assistant"]


def test_chinese_font_fallback_reaches_shapes_and_speaker_notes(tmp_path: Path) -> None:
    """shape 与 speaker notes 都只替换中文使用的 a:ea slot。"""

    source = tmp_path / "notes-source.pptx"
    output = tmp_path / "notes-translated.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    body_shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    body_shape.text = "Body"
    notes_text_frame = slide.notes_slide.notes_text_frame
    assert notes_text_frame is not None
    notes_text_frame.text = "Note"
    source_body_run = body_shape.text_frame.paragraphs[0].runs[0]
    source_notes_run = notes_text_frame.paragraphs[0].runs[0]
    for run in (source_body_run, source_notes_run):
        _set_run_font_slots(
            run,
            latin="Aptos",
            east_asian="Wingdings",
            complex_script="Nirmala UI",
        )
    presentation.save(source)

    translator = RecordingTranslator(
        lambda text: text.replace("Body", "正文").replace("Note", "演讲者备注")
    )
    PptxPipeline(translator, translate_tables=False).translate_to(
        source_path=source,
        output_path=output,
        source_language="en",
        target_language="zh-CN",
    )

    translated_slide = Presentation(output).slides[0]
    body_run = translated_slide.shapes[0].text_frame.paragraphs[0].runs[0]
    notes_run = translated_slide.notes_slide.notes_text_frame.paragraphs[0].runs[0]
    for run in (body_run, notes_run):
        run_properties = run._r.get_or_add_rPr()
        assert run_properties.find(qn("a:latin")).get("typeface") == "Aptos"
        assert run_properties.find(qn("a:ea")).get("typeface") == "Noto Sans SC"
        assert run_properties.find(qn("a:cs")).get("typeface") == "Nirmala UI"


def test_east_asian_compatible_source_font_is_preserved(tmp_path: Path) -> None:
    """PPTX 未声明 a:ea 时仍按 a:latin 识别兼容中文字体。"""
    source = tmp_path / "east-asian-font.pptx"
    output = tmp_path / "translated.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    paragraph = slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(5), Inches(1)
    ).text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = "Original"
    run.font.name = "Songti SC"
    presentation.save(source)

    PptxPipeline(
        RecordingTranslator(lambda text: text.replace("Original", "中文")),
        translate_tables=False,
        translate_notes=False,
    ).translate_to(
        source_path=source,
        output_path=output,
        source_language="en",
        target_language="zh-CN",
    )

    translated_run = Presentation(output).slides[0].shapes[0].text_frame.paragraphs[0].runs[0]
    translated_properties = translated_run._r.get_or_add_rPr()
    assert translated_properties.find(qn("a:ea")) is None
    assert translated_properties.find(qn("a:latin")).get("typeface") == "Songti SC"


def test_latin_fallback_only_replaces_latin_font_slot(tmp_path: Path) -> None:
    """English ASCII 遇到 Wingdings 时只替换 a:latin。"""

    source = tmp_path / "latin-font.pptx"
    output = tmp_path / "translated.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    paragraph = slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(5), Inches(1)
    ).text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = "原文"
    _set_run_font_slots(
        run,
        latin="Wingdings",
        east_asian="Songti SC",
        complex_script="Nirmala UI",
    )
    presentation.save(source)

    PptxPipeline(
        RecordingTranslator(lambda text: text.replace("原文", "Contract terms")),
        translate_tables=False,
        translate_notes=False,
    ).translate_to(
        source_path=source,
        output_path=output,
        source_language="zh-CN",
        target_language="en",
    )

    translated_run = Presentation(output).slides[0].shapes[0].text_frame.paragraphs[0].runs[0]
    run_properties = translated_run._r.get_or_add_rPr()
    assert run_properties.find(qn("a:latin")).get("typeface") == "Calibri"
    assert run_properties.find(qn("a:ea")).get("typeface") == "Songti SC"
    assert run_properties.find(qn("a:cs")).get("typeface") == "Nirmala UI"


def test_malformed_translation_gets_one_repair_then_falls_back(tmp_path: Path) -> None:
    """损坏的 span 只 repair 一次, 仍失败时保留原段落并返回 warning."""

    source = _single_text_presentation(tmp_path / "source.pptx", "Original")
    output = tmp_path / "translated.pptx"
    translator = RecordingTranslator(lambda _text: "<span><SPACE></span>")

    result = PptxPipeline(translator, translate_tables=False, translate_notes=False).translate_to(
        source_path=source,
        output_path=output,
        source_language="en",
        target_language="fr",
    )

    assert Presentation(output).slides[0].shapes[0].text == "Original"
    assert [hint for hint, _ in translator.calls] == ["pptx", "pptx_repair"]
    assert result.translated_segments == 0
    assert result.fallback_segments == 1
    assert result.warning_codes == ("pptx_segment_fallback",)


def test_nested_span_translation_uses_joto_deterministic_repair(tmp_path: Path) -> None:
    """模型误嵌套 span 时按兼容规则恢复三个 run。"""

    source = tmp_path / "nested-span-source.pptx"
    output = tmp_path / "nested-span-translated.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    paragraph = slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(7), Inches(1)
    ).text_frame.paragraphs[0]
    for index, source_text in enumerate(("前半句", "关键词", "后半句")):
        run = paragraph.add_run()
        run.text = source_text
        run.font.bold = index == 1
        run.font.italic = index == 2
    presentation.save(source)

    malformed = "<span>Text before <span>keyword</span> text after</span>"
    translator = RecordingTranslator(lambda _text: malformed)
    result = PptxPipeline(translator, translate_tables=False, translate_notes=False).translate_to(
        source_path=source,
        output_path=output,
        source_language="zh-CN",
        target_language="en",
    )

    translated_runs = Presentation(output).slides[0].shapes[0].text_frame.paragraphs[0].runs
    assert [run.text for run in translated_runs] == ["Text before ", "keyword", " text after"]
    assert [hint for hint, _ in translator.calls] == ["pptx", "pptx_repair"]
    assert result.translated_segments == 1
    assert result.fallback_segments == 0


def test_malformed_segments_share_one_batch_repair_call(tmp_path: Path) -> None:
    """同一 group 的多个坏 segment 只触发一次 batch repair。"""

    source = tmp_path / "batch-repair-source.pptx"
    output = tmp_path / "batch-repair-translated.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    for top, text in ((1, "Alpha"), (2, "Beta")):
        paragraph = slide.shapes.add_textbox(
            Inches(1), Inches(top), Inches(5), Inches(1)
        ).text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = text
    presentation.save(source)

    translator = RecordingTranslator(lambda _text: "<span>translated")
    result = PptxPipeline(translator, translate_tables=False, translate_notes=False).translate_to(
        source_path=source,
        output_path=output,
        source_language="en",
        target_language="fr",
    )

    assert [hint for hint, _ in translator.calls] == ["pptx", "pptx_repair"]
    assert len(translator.calls[1][1]) == 2
    translated_slide = Presentation(output).slides[0]
    assert [shape.text for shape in translated_slide.shapes] == ["translated", "translated"]
    assert result.translated_segments == 2
    assert result.fallback_segments == 0


def test_provider_failure_falls_back_without_failing_the_file(tmp_path: Path) -> None:
    """provider 整批失败时直接 fallback, 不产生逐段 repair 调用。"""

    source = _single_text_presentation(tmp_path / "source.pptx", "Original")
    output = tmp_path / "translated.pptx"
    translator = FailingTranslator()

    result = PptxPipeline(translator, translate_tables=False, translate_notes=False).translate_to(
        source_path=source,
        output_path=output,
        source_language=None,
        target_language="fr",
    )

    assert Presentation(output).slides[0].shapes[0].text == "Original"
    assert translator.format_hints == ["pptx"]
    assert result.fallback_segments == 1


def test_missing_candidate_falls_back_without_repair(tmp_path: Path) -> None:
    """provider 遗漏 candidate 时直接保留原文, 不凭空发起 repair。"""

    source = _single_text_presentation(tmp_path / "source.pptx", "Original")
    output = tmp_path / "translated.pptx"
    translator = MissingCandidateTranslator()

    result = PptxPipeline(translator, translate_tables=False, translate_notes=False).translate_to(
        source_path=source,
        output_path=output,
        source_language="en",
        target_language="fr",
    )

    assert Presentation(output).slides[0].shapes[0].text == "Original"
    assert translator.format_hints == ["pptx"]
    assert result.fallback_segments == 1


@pytest.mark.parametrize("boundary_tag", ["a:br", "a:fld"])
def test_run_normalizer_respects_non_run_xml_boundaries(boundary_tag: str) -> None:
    """视觉相同的 run 也不得跨 a:br 或 a:fld 合并。"""

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    paragraph = slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(5), Inches(1)
    ).text_frame.paragraphs[0]
    first_run = paragraph.add_run()
    first_run.text = "Alpha"
    second_run = paragraph.add_run()
    second_run.text = "Beta"
    boundary = OxmlElement(boundary_tag)
    first_run._r.addnext(boundary)

    normalized = PptxRunNormalizer().merge_runs(list(paragraph.runs))

    assert [run.text for run in normalized] == ["Alpha", "Beta"]
    assert [run.source_run_indices for run in normalized] == [(0,), (1,)]


def test_run_normalizer_includes_hyperlink_attributes_in_merge_key() -> None:
    """相同 rId 但 action 或 tooltip 不同的 run 必须保留独立语义边界。"""

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    paragraph = slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(5), Inches(1)
    ).text_frame.paragraphs[0]
    first_run = paragraph.add_run()
    first_run.text = "Alpha"
    first_run.hyperlink.address = "https://example.com/shared"
    second_run = paragraph.add_run()
    second_run.text = "Beta"
    second_run.hyperlink.address = "https://example.com/shared"
    first_hyperlink = _click_hyperlink(first_run)
    second_hyperlink = _click_hyperlink(second_run)
    second_hyperlink.set(qn("r:id"), first_hyperlink.get(qn("r:id")))
    first_hyperlink.set("action", "ppaction://hlinkshowjump?jump=firstslide")
    first_hyperlink.set("tooltip", "first target")
    second_hyperlink.set("action", "ppaction://hlinkshowjump?jump=lastslide")
    second_hyperlink.set("tooltip", "last target")

    normalized = PptxRunNormalizer().merge_runs(list(paragraph.runs))

    assert [run.text for run in normalized] == ["Alpha", "Beta"]
    assert (
        normalized[0].format_info["hyperlink_semantics"]
        != normalized[1].format_info["hyperlink_semantics"]
    )


def test_fallback_preserves_line_break_and_two_hyperlinks(tmp_path: Path) -> None:
    """fallback 后仍保留换行位置、双链接文字及其完整 hyperlink 属性。"""

    source = _linked_boundary_presentation(tmp_path / "source.pptx")
    output = tmp_path / "translated.pptx"
    source_paragraph = Presentation(source).slides[0].shapes[0].text_frame.paragraphs[0]
    expected_semantics = _paragraph_semantics(source_paragraph)
    translator = FailingTranslator()

    result = PptxPipeline(translator, translate_tables=False, translate_notes=False).translate_to(
        source_path=source,
        output_path=output,
        source_language="en",
        target_language="fr",
    )

    output_paragraph = Presentation(output).slides[0].shapes[0].text_frame.paragraphs[0]
    assert _paragraph_semantics(output_paragraph) == expected_semantics
    assert [run.text for run in output_paragraph.runs] == ["Alpha", "Beta", "Gamma"]
    assert _visible_child_tags(output_paragraph) == ("r", "r", "br", "r")
    assert translator.format_hints == ["pptx"]
    assert result.fallback_segments == 1


def test_validation_failure_keeps_existing_output_and_cleans_temp_file(
    tmp_path: Path,
) -> None:
    """保存后的 package 校验失败时不修改现有目标文件."""

    source = _single_text_presentation(tmp_path / "source.pptx", "Original")
    output = tmp_path / "translated.pptx"
    output.write_bytes(b"existing-output")
    translator = RecordingTranslator(lambda text: text.replace("Original", "Translated"))

    class InvalidatingPipeline(PptxPipeline):
        """强制在原子替换前触发失败."""

        def _validate_output(self, temporary_path, expected_signature) -> None:
            """模拟结构校验拒绝临时 package."""

            del temporary_path, expected_signature
            raise ValueError("synthetic validation failure")

    with pytest.raises(ValueError, match="synthetic validation failure"):
        InvalidatingPipeline(
            translator, translate_tables=False, translate_notes=False
        ).translate_to(
            source_path=source,
            output_path=output,
            source_language="en",
            target_language="fr",
        )

    assert output.read_bytes() == b"existing-output"
    assert not list(tmp_path.glob(".translated.*.tmp.pptx"))


def _set_run_font_slots(
    run,
    *,
    latin: str,
    east_asian: str,
    complex_script: str,
) -> None:
    """给测试 run 写入三个独立 DrawingML 字体 slot。"""

    run.font.name = latin
    run_properties = run._r.get_or_add_rPr()
    for tag, font_name in (("a:ea", east_asian), ("a:cs", complex_script)):
        typeface = OxmlElement(tag)
        typeface.set("typeface", font_name)
        run_properties.append(typeface)


def _single_text_presentation(path: Path, text: str) -> Path:
    """创建只有一个 slide 和一个文本框的 presentation."""

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text = text
    presentation.save(path)
    return path


def _linked_boundary_presentation(path: Path) -> Path:
    """创建含相邻双 hyperlink 和后置 line break 的测试 presentation。"""

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    paragraph = slide.shapes.add_textbox(
        Inches(1), Inches(1), Inches(5), Inches(2)
    ).text_frame.paragraphs[0]
    first_run = paragraph.add_run()
    first_run.text = "Alpha"
    first_run.hyperlink.address = "https://example.com/alpha"
    first_hyperlink = _click_hyperlink(first_run)
    first_hyperlink.set("action", "ppaction://hlinkshowjump?jump=firstslide")
    first_hyperlink.set("tooltip", "Alpha tooltip")

    second_run = paragraph.add_run()
    second_run.text = "Beta"
    second_run.hyperlink.address = "https://example.com/beta"
    second_hyperlink = _click_hyperlink(second_run)
    second_hyperlink.set("action", "ppaction://hlinkshowjump?jump=lastslide")
    second_hyperlink.set("tooltip", "Beta tooltip")

    paragraph.add_line_break()
    final_run = paragraph.add_run()
    final_run.text = "Gamma"
    presentation.save(path)
    return path


def _click_hyperlink(run):
    """返回测试 run 已存在的 a:hlinkClick 节点。"""

    run_properties = run._r.get_or_add_rPr()
    hyperlink = run_properties.find(qn("a:hlinkClick"))
    assert hyperlink is not None
    return hyperlink


def _visible_child_tags(paragraph) -> tuple[str, ...]:
    """返回段落中参与文字顺序的 DrawingML 子节点名称。"""

    return tuple(
        child.tag.split("}", 1)[-1]
        for child in paragraph._p
        if child.tag.split("}", 1)[-1] in {"r", "br", "fld"}
    )


def _paragraph_semantics(paragraph) -> tuple[tuple[str, ...], tuple[tuple, ...]]:
    """记录 run 文本、hyperlink 属性和 relationship target。"""

    run_semantics: list[tuple] = []
    for run in paragraph.runs:
        hyperlink = _click_hyperlink(run) if run.hyperlink.address else None
        if hyperlink is None:
            run_semantics.append((run.text, (), None))
            continue
        relationship_id = hyperlink.get(qn("r:id"))
        target = run.part.rels[relationship_id].target_ref if relationship_id else None
        run_semantics.append(
            (
                run.text,
                tuple(sorted(hyperlink.attrib.items())),
                target,
            )
        )
    return _visible_child_tags(paragraph), tuple(run_semantics)


def _notes_relationship(slide) -> tuple[str, str]:
    """返回 slide 到 notes 的 relationship 标识."""

    relationships = [
        relationship
        for relationship in slide.part.rels.values()
        if relationship.reltype.endswith("/notesSlide")
    ]
    assert len(relationships) == 1
    relationship = relationships[0]
    return relationship.rId, str(relationship.target_part.partname)


def _non_body_notes_placeholders(slide) -> tuple[tuple[int, str], ...]:
    """描述翻译路径绝不能修改的非 BODY notes placeholder."""

    notes_slide = slide.notes_slide
    values: list[tuple[int, str]] = []
    for placeholder in notes_slide.placeholders:
        if placeholder.placeholder_format.type == PP_PLACEHOLDER.BODY:
            continue
        text = placeholder.text if getattr(placeholder, "has_text_frame", False) else ""
        values.append((int(placeholder.placeholder_format.type), text))
    return tuple(values)
