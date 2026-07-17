"""验证 Office run 合并不会跨越不同的字体 slot 语义边界。"""

from typing import Any

import pytest
from docx import Document
from docx.oxml import OxmlElement as DocxOxmlElement
from docx.oxml.ns import qn as docx_qn
from pptx import Presentation
from pptx.oxml.ns import qn as pptx_qn
from pptx.oxml.xmlchemy import OxmlElement as PptxOxmlElement
from pptx.util import Inches

from modules.docx.extractor import DocxExtractor
from modules.pptx.run_normalizer import PptxRunNormalizer


def _set_word_font_slots(run: Any, attributes: dict[str, str]) -> None:
    """给 DOCX run 写入一组显式 ``w:rFonts`` 属性。"""

    run_properties = run._element.get_or_add_rPr()
    run_fonts = run_properties.find(docx_qn("w:rFonts"))
    if run_fonts is None:
        run_fonts = DocxOxmlElement("w:rFonts")
        run_properties.insert(0, run_fonts)
    for name, value in attributes.items():
        run_fonts.set(docx_qn(f"w:{name}"), value)


def _new_pptx_paragraph() -> Any:
    """创建含普通 text frame 的 PPTX 段落供边界测试复用。"""

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    return slide.shapes.add_textbox(
        Inches(1),
        Inches(1),
        Inches(5),
        Inches(1),
    ).text_frame.paragraphs[0]


def _set_drawingml_font_slot(
    run: Any,
    slot: str,
    *,
    typeface: str,
    extra_attributes: dict[str, str] | None = None,
) -> None:
    """给 PPTX run 写入一个完整 DrawingML font slot。"""

    run_properties = run._r.get_or_add_rPr()
    font_slot = run_properties.find(pptx_qn(f"a:{slot}"))
    if font_slot is None:
        font_slot = PptxOxmlElement(f"a:{slot}")
        run_properties.append(font_slot)
    font_slot.set("typeface", typeface)
    for name, value in (extra_attributes or {}).items():
        font_slot.set(name, value)


@pytest.mark.parametrize(
    ("slot", "first_value", "second_value"),
    [
        ("ascii", "Arial", "Calibri"),
        ("hAnsi", "Arial", "Calibri"),
        ("eastAsia", "SimSun", "Microsoft YaHei"),
        ("cs", "Arial", "Times New Roman"),
        ("asciiTheme", "majorAscii", "minorAscii"),
        ("hAnsiTheme", "majorHAnsi", "minorHAnsi"),
        ("eastAsiaTheme", "majorEastAsia", "minorEastAsia"),
        ("cstheme", "majorBidi", "minorBidi"),
        ("hint", "default", "eastAsia"),
    ],
)
def test_docx_merge_respects_every_word_font_slot(
    slot: str,
    first_value: str,
    second_value: str,
) -> None:
    """任一 ``w:rFonts`` 属性不同都必须保留两个 DOCX span。"""

    document = Document()
    paragraph = document.add_paragraph()
    first_run = paragraph.add_run("Alpha")
    second_run = paragraph.add_run("Beta")
    common_slots = {"ascii": "Arial", "hAnsi": "Arial"}
    _set_word_font_slots(first_run, common_slots | {slot: first_value})
    _set_word_font_slots(second_run, common_slots | {slot: second_value})

    normalized = DocxExtractor().extract_runs(paragraph)

    assert [run.text for run in normalized] == ["Alpha", "Beta"]
    assert normalized[0].format_info["font_slots"] != normalized[1].format_info["font_slots"]


def test_docx_token_projection_keeps_original_font_merge_behavior() -> None:
    """真实 span 按字体拆分时, token projection 仍维持原 batch 文本。"""

    document = Document()
    paragraph = document.add_paragraph()
    first_run = paragraph.add_run("Alpha")
    second_run = paragraph.add_run("Beta")
    _set_word_font_slots(
        first_run,
        {"ascii": "Arial", "hAnsi": "Arial", "eastAsia": "SimSun"},
    )
    _set_word_font_slots(
        second_run,
        {"ascii": "Arial", "hAnsi": "Arial", "eastAsia": "Microsoft YaHei"},
    )
    extractor = DocxExtractor()

    normalized = extractor.extract_runs(paragraph)
    token_projection = extractor._build_batch_token_projection(paragraph, 0)

    assert [run.text for run in normalized] == ["Alpha", "Beta"]
    assert token_projection == "[PARA_0]<span>AlphaBeta</span>"


@pytest.mark.parametrize("slot", ["latin", "ea", "cs"])
def test_pptx_merge_respects_all_drawingml_font_slots(slot: str) -> None:
    """``a:latin``、``a:ea`` 或 ``a:cs`` 不同都不得合并 PPTX run。"""

    paragraph = _new_pptx_paragraph()
    first_run = paragraph.add_run()
    first_run.text = "Alpha"
    first_run.font.name = "Arial"
    second_run = paragraph.add_run()
    second_run.text = "Beta"
    second_run.font.name = "Arial"
    _set_drawingml_font_slot(first_run, slot, typeface="First Font")
    _set_drawingml_font_slot(second_run, slot, typeface="Second Font")

    normalized = PptxRunNormalizer().merge_runs(list(paragraph.runs))

    assert [run.text for run in normalized] == ["Alpha", "Beta"]
    assert normalized[0].format_info["font_slots"] != normalized[1].format_info["font_slots"]


def test_pptx_merge_respects_font_slot_metadata() -> None:
    """typeface 相同但 charset 不同的 DrawingML font slot 也保持独立。"""

    paragraph = _new_pptx_paragraph()
    first_run = paragraph.add_run()
    first_run.text = "Alpha"
    second_run = paragraph.add_run()
    second_run.text = "Beta"
    _set_drawingml_font_slot(
        first_run,
        "ea",
        typeface="Arial Unicode MS",
        extra_attributes={"charset": "01"},
    )
    _set_drawingml_font_slot(
        second_run,
        "ea",
        typeface="Arial Unicode MS",
        extra_attributes={"charset": "02"},
    )

    normalized = PptxRunNormalizer().merge_runs(list(paragraph.runs))

    assert [run.text for run in normalized] == ["Alpha", "Beta"]
    assert normalized[0].format_info["font_slots"] != normalized[1].format_info["font_slots"]
